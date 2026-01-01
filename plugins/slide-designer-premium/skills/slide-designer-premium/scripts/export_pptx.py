#!/usr/bin/env python3
"""
Export HTML slide deck to PPTX format.
Parses HTML slides and recreates in PowerPoint with styling.

Usage:
    python3 export_pptx.py input.html output.pptx [--theme orchestrate]
"""

import sys
import re
from pathlib import Path
from bs4 import BeautifulSoup

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
except ImportError:
    print("Error: python-pptx not installed. Run: pip install python-pptx")
    sys.exit(1)

try:
    from bs4 import BeautifulSoup
except ImportError:
    print("Error: beautifulsoup4 not installed. Run: pip install beautifulsoup4")
    sys.exit(1)


# Theme color mappings
THEMES = {
    "orchestrate": {
        "bg": "0A0A0F",
        "accent": "F97316",
        "text": "FFFFFF",
        "text_secondary": "A1A1AA"
    },
    "corporate-blue": {
        "bg": "0F172A",
        "accent": "3B82F6",
        "text": "F8FAFC",
        "text_secondary": "94A3B8"
    },
    "minimal-light": {
        "bg": "FFFFFF",
        "accent": "0F172A",
        "text": "0F172A",
        "text_secondary": "475569"
    },
    "startup-gradient": {
        "bg": "0D0D12",
        "accent": "8B5CF6",
        "text": "FFFFFF",
        "text_secondary": "A1A1AA"
    },
    "tome": {
        "bg": "000000",
        "accent": "8B5CF6",
        "text": "FFFFFF",
        "text_secondary": "A0A0A0"
    },
    "gamma": {
        "bg": "FFFFFF",
        "accent": "7C3AED",
        "text": "18181B",
        "text_secondary": "71717A"
    },
    "beautiful-ai": {
        "bg": "FFFFFF",
        "accent": "3B82F6",
        "text": "0F172A",
        "text_secondary": "64748B"
    },
    "beautiful-ai-dark": {
        "bg": "0F172A",
        "accent": "3B82F6",
        "text": "F8FAFC",
        "text_secondary": "94A3B8"
    }
}


def get_text_content(element):
    """Extract all text content, preserving spaces between nested elements."""
    if element is None:
        return ""
    # Use separator to keep spaces between elements like <strong>
    text = element.get_text(separator=' ', strip=True)
    # Clean up multiple spaces
    import re
    return re.sub(r'\s+', ' ', text).strip()


def parse_html_slides(html_content):
    """Parse HTML content and extract slide data using BeautifulSoup."""
    soup = BeautifulSoup(html_content, 'html.parser')
    slides = []

    for slide_div in soup.find_all('div', class_='slide'):
        slide_data = {
            'title': '',
            'subtitle': '',
            'overline': '',
            'bullets': [],
            'content_blocks': [],
            'type': 'content'
        }

        # Detect slide type
        classes = slide_div.get('class', [])
        if 'title-slide' in classes:
            slide_data['type'] = 'title'
        elif 'section-divider' in classes:
            slide_data['type'] = 'section'

        # Extract overline
        overline = slide_div.find(class_='overline')
        if overline:
            slide_data['overline'] = get_text_content(overline)

        # Extract title (h1 or h2)
        h1 = slide_div.find('h1')
        h2 = slide_div.find('h2')
        if h1:
            slide_data['title'] = get_text_content(h1)
        elif h2:
            slide_data['title'] = get_text_content(h2)

        # Extract subtitle
        subtitle = slide_div.find(class_='subtitle')
        if subtitle:
            slide_data['subtitle'] = get_text_content(subtitle)

        # Extract bullets from ul/ol lists
        for ul in slide_div.find_all(['ul', 'ol']):
            for li in ul.find_all('li', recursive=False):
                text = get_text_content(li)
                if text:
                    slide_data['bullets'].append(text)

        # Extract timeline items
        for timeline_item in slide_div.find_all(class_='timeline-item'):
            date = timeline_item.find(class_='timeline-date')
            title = timeline_item.find(class_='timeline-title')
            desc = timeline_item.find(class_='timeline-desc')

            parts = []
            if date:
                parts.append(get_text_content(date))
            if title:
                parts.append(get_text_content(title))
            if desc:
                parts.append(get_text_content(desc))

            if parts:
                slide_data['content_blocks'].append(' - '.join(parts))

        # Extract feature cards
        for card in slide_div.find_all(class_='feature-card'):
            card_title = card.find('h3') or card.find(class_='card-title')
            card_desc = card.find('p') or card.find(class_='card-desc')

            parts = []
            if card_title:
                parts.append(get_text_content(card_title))
            if card_desc:
                parts.append(get_text_content(card_desc))

            if parts:
                slide_data['content_blocks'].append(' - '.join(parts))

        # Extract stat items (class='stat' or 'stat-card')
        for stat in slide_div.find_all(class_=lambda c: c and ('stat' in c or 'stat-card' in c)):
            value = stat.find(class_='stat-value')
            label = stat.find(class_='stat-label')
            context = stat.find(class_='stat-context')

            parts = []
            if value:
                parts.append(get_text_content(value))
            if label:
                parts.append(get_text_content(label))
            if context:
                parts.append(f"({get_text_content(context)})")

            if parts:
                slide_data['content_blocks'].append(' '.join(parts))

        # Extract action items
        for action in slide_div.find_all(class_='action-item'):
            text = get_text_content(action)
            if text:
                slide_data['content_blocks'].append(text)

        # Extract decision items (class='decision-item' or 'decision-card')
        for decision in slide_div.find_all(class_=lambda c: c and ('decision-item' in c or 'decision-card' in c)):
            title = decision.find(class_='decision-title')
            detail = decision.find(class_='decision-detail')

            parts = []
            if title:
                parts.append(get_text_content(title))
            if detail:
                parts.append(get_text_content(detail))

            if parts:
                slide_data['content_blocks'].append(' - '.join(parts))
            elif get_text_content(decision):
                slide_data['content_blocks'].append(get_text_content(decision))

        # Extract table data
        for table in slide_div.find_all('table'):
            for row in table.find_all('tr'):
                cells = row.find_all(['td', 'th'])
                if cells:
                    row_text = ' | '.join(get_text_content(cell) for cell in cells)
                    if row_text.strip():
                        slide_data['content_blocks'].append(row_text)

        # Extract column lists (two-column layouts)
        for column in slide_div.find_all(class_='column'):
            header = column.find(class_='column-header')
            if header:
                slide_data['content_blocks'].append(get_text_content(header))
            for li in column.find_all('li'):
                text = get_text_content(li)
                if text and text not in slide_data['bullets']:
                    slide_data['bullets'].append(text)

        # Extract region cards
        for region in slide_div.find_all(class_='region-card'):
            name = region.find(class_='region-name')
            percent = region.find(class_='region-percent')
            revenue = region.find(class_='region-revenue')

            parts = []
            if name:
                parts.append(get_text_content(name))
            if percent:
                parts.append(get_text_content(percent))
            if revenue:
                parts.append(get_text_content(revenue))

            if parts:
                slide_data['content_blocks'].append(' - '.join(parts))

        # Extract insight cards
        for insight in slide_div.find_all(class_='insight-card'):
            title = insight.find('h3')
            desc = insight.find('p')

            parts = []
            if title:
                parts.append(get_text_content(title))
            if desc:
                parts.append(get_text_content(desc))

            if parts:
                slide_data['content_blocks'].append(' - '.join(parts))

        # Extract meeting info (next-meeting-card)
        meeting_card = slide_div.find(class_='next-meeting-card')
        if meeting_card:
            date = meeting_card.find(class_='meeting-date')
            time = meeting_card.find(class_='meeting-time')
            if date or time:
                parts = ['Next Meeting:']
                if date:
                    parts.append(get_text_content(date))
                if time:
                    parts.append(get_text_content(time))
                slide_data['content_blocks'].append(' '.join(parts))

            # Agenda items
            for li in meeting_card.find_all('li'):
                text = get_text_content(li)
                if text:
                    slide_data['bullets'].append(text)

        # Extract any remaining paragraphs not in special containers
        for p in slide_div.find_all('p', recursive=False):
            if 'subtitle' not in p.get('class', []):
                text = get_text_content(p)
                if text and text not in slide_data['bullets']:
                    slide_data['content_blocks'].append(text)

        slides.append(slide_data)

    return slides


def hex_to_rgb(hex_color):
    """Convert hex color to RGB tuple."""
    hex_color = hex_color.lstrip('#')
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))


def create_pptx(slides, output_path, theme_name="orchestrate"):
    """Create PPTX from parsed slide data."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9
    prs.slide_height = Inches(7.5)

    theme = THEMES.get(theme_name, THEMES["orchestrate"])
    bg_rgb = hex_to_rgb(theme["bg"])
    accent_rgb = hex_to_rgb(theme["accent"])
    text_rgb = hex_to_rgb(theme["text"])
    text_secondary_rgb = hex_to_rgb(theme["text_secondary"])

    for slide_data in slides:
        # Use blank layout
        slide_layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(slide_layout)

        # Set background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*bg_rgb)

        # Positioning
        left_margin = Inches(0.75)
        top_margin = Inches(0.8)

        if slide_data['type'] == 'title':
            # Title slide - centered
            # Overline
            if slide_data.get('overline'):
                overline_box = slide.shapes.add_textbox(
                    Inches(0.75), Inches(2.5), Inches(11.833), Inches(0.5)
                )
                overline_frame = overline_box.text_frame
                overline_para = overline_frame.paragraphs[0]
                overline_para.text = slide_data['overline'].upper()
                overline_para.font.size = Pt(14)
                overline_para.font.color.rgb = RGBColor(*accent_rgb)
                overline_para.font.bold = True
                overline_para.alignment = PP_ALIGN.CENTER

            # Title
            title_box = slide.shapes.add_textbox(
                Inches(0.75), Inches(3), Inches(11.833), Inches(1.5)
            )
            title_frame = title_box.text_frame
            title_para = title_frame.paragraphs[0]
            title_para.text = slide_data['title']
            title_para.font.size = Pt(56)
            title_para.font.color.rgb = RGBColor(*text_rgb)
            title_para.font.bold = True
            title_para.alignment = PP_ALIGN.CENTER

            # Subtitle
            if slide_data.get('subtitle'):
                subtitle_box = slide.shapes.add_textbox(
                    Inches(0.75), Inches(4.5), Inches(11.833), Inches(0.75)
                )
                subtitle_frame = subtitle_box.text_frame
                subtitle_para = subtitle_frame.paragraphs[0]
                subtitle_para.text = slide_data['subtitle']
                subtitle_para.font.size = Pt(24)
                subtitle_para.font.color.rgb = RGBColor(*text_secondary_rgb)
                subtitle_para.alignment = PP_ALIGN.CENTER

        else:
            # Content slide
            # Title
            title_box = slide.shapes.add_textbox(
                left_margin, top_margin, Inches(11.833), Inches(1)
            )
            title_frame = title_box.text_frame
            title_para = title_frame.paragraphs[0]
            title_para.text = slide_data['title']
            title_para.font.size = Pt(42)
            title_para.font.color.rgb = RGBColor(*text_rgb)
            title_para.font.bold = True

            # Accent bar
            accent_bar = slide.shapes.add_shape(
                1,  # Rectangle
                left_margin, Inches(1.9), Inches(1), Pt(4)
            )
            accent_bar.fill.solid()
            accent_bar.fill.fore_color.rgb = RGBColor(*accent_rgb)
            accent_bar.line.fill.background()

            # Combine bullets and content_blocks
            all_content = slide_data['bullets'] + slide_data['content_blocks']

            if all_content:
                content_box = slide.shapes.add_textbox(
                    left_margin, Inches(2.3), Inches(11.833), Inches(4.5)
                )
                content_frame = content_box.text_frame
                content_frame.word_wrap = True

                for i, item in enumerate(all_content):
                    if i == 0:
                        para = content_frame.paragraphs[0]
                    else:
                        para = content_frame.add_paragraph()
                    para.text = f"• {item}"
                    para.font.size = Pt(22)
                    para.font.color.rgb = RGBColor(*text_rgb)
                    para.space_after = Pt(14)

    prs.save(output_path)
    return output_path


def main():
    if len(sys.argv) < 3:
        print("Usage: python3 export_pptx.py input.html output.pptx [--theme orchestrate]")
        print("\nAvailable themes: orchestrate, corporate-blue, minimal-light, startup-gradient, tome, gamma, beautiful-ai, beautiful-ai-dark")
        sys.exit(1)

    input_path = sys.argv[1]
    output_path = sys.argv[2]
    theme = "orchestrate"

    # Parse optional theme argument
    if "--theme" in sys.argv:
        theme_idx = sys.argv.index("--theme")
        if theme_idx + 1 < len(sys.argv):
            theme = sys.argv[theme_idx + 1]

    # Read HTML
    with open(input_path, 'r') as f:
        html_content = f.read()

    # Parse slides
    slides = parse_html_slides(html_content)

    if not slides:
        print("ERROR: No slides found in HTML. Check that slides have class='slide'")
        sys.exit(1)

    print(f"Found {len(slides)} slides")
    for i, s in enumerate(slides):
        bullet_count = len(s['bullets']) + len(s['content_blocks'])
        print(f"  Slide {i+1}: {s['title'][:40]}... ({bullet_count} items)")

    # Create PPTX
    create_pptx(slides, output_path, theme)
    print(f"Created: {output_path}")


if __name__ == "__main__":
    main()
